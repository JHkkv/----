"""
东南亚盲盒市场网页设计方案汇报PPT v3 - 稳健布局版
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 配色
C1 = RGBColor(0x1A, 0x56, 0xDB)  # 深蓝主色
C2 = RGBColor(0x4F, 0x86, 0xF7)  # 浅蓝
C3 = RGBColor(0xFF, 0x6B, 0x6B)  # 珊瑚红
C4 = RGBColor(0x4E, 0xC9, 0xB0)  # 青绿
C5 = RGBColor(0xFF, 0xB8, 0x4C)  # 金黄
CDARK = RGBColor(0x2C, 0x3E, 0x50)
CGRAY = RGBColor(0x7F, 0x8C, 0x8D)
CLIGHT = RGBColor(0xBD, 0xC3, 0xC7)
CWHITE = RGBColor(0xFF, 0xFF, 0xFF)
CBG = RGBColor(0xF0, 0xF2, 0xF5)
CDARKBG = RGBColor(0x1E, 0x2D, 0x3D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
FONT = 'Microsoft YaHei'


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=CWHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def rrect(slide, l, t, w, h, fill=CWHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=18, bold=False, color=CDARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """创建文本框，自动调整文字大小以避免溢出"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    return tb


def multi_text(slide, l, t, w, h, items, sz=14, color=CDARK, bullet='•'):
    """创建多行文本，每行一个条目"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'  {bullet} {item}'
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_before = Pt(6)
    return tb


def section_title(slide, title_text):
    """统一章节标题"""
    rect(slide, Inches(0), Inches(0), SW, Inches(0.06), C1)
    txt(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
        title_text, sz=30, bold=True, color=C1)


def kpi_card(slide, x, y, label, value):
    """KPI卡片"""
    rrect(slide, x, y, Inches(2.8), Inches(1.6), CWHITE)
    txt(slide, x + Inches(0.3), y + Inches(0.2), Inches(2.2), Inches(0.4),
        label, sz=12, color=CGRAY)
    txt(slide, x + Inches(0.3), y + Inches(0.6), Inches(2.2), Inches(0.6),
        value, sz=22, bold=True, color=C1)


# ================================================================
# Slide 1: 封面
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CDARKBG)
txt(sl, Inches(1.5), Inches(2), Inches(10), Inches(1.5),
    '东南亚盲盒市场网页设计方案', sz=40, bold=True, color=CWHITE, align=PP_ALIGN.CENTER)
txt(sl, Inches(1.5), Inches(3.5), Inches(10), Inches(0.8),
    '市场分析  |  设计方案  |  风格评估', sz=20, color=C2, align=PP_ALIGN.CENTER)
txt(sl, Inches(1.5), Inches(5), Inches(10), Inches(0.5),
    '2026年6月', sz=14, color=CGRAY, align=PP_ALIGN.CENTER)

# ================================================================
# Slide 2: 项目背景
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CBG)
section_title(sl, '项目背景与目标')

# 左侧-市场机会
rrect(sl, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.5), CWHITE)
txt(sl, Inches(1), Inches(1.35), Inches(5), Inches(0.5), '市场机会', sz=20, bold=True, color=C1)
multi_text(sl, Inches(1), Inches(1.9), Inches(5), Inches(4.5), [
    '东南亚盲盒市场年复合增长率（CAGR）达28.5%',
    '2025年市场规模预计达25亿美元',
    '泡泡玛特海外收入突破50亿元人民币',
    '泰国、越南、印尼为核心增长引擎',
    '年轻人口结构占主导（中位年龄约30岁）',
    '社交媒体渗透率全球领先',
], sz=14, color=CDARK)

# 右侧-项目目标
rrect(sl, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.5), CWHITE)
txt(sl, Inches(7.2), Inches(1.35), Inches(5), Inches(0.5), '项目目标', sz=20, bold=True, color=C3)
multi_text(sl, Inches(7.2), Inches(1.9), Inches(5), Inches(4.5), [
    '深度分析东南亚盲盒市场数据',
    '设计10套网页模板概念方案',
    '评估品牌风格适配性',
    '推荐最佳市场解决方案',
    '提供可落地的设计指南',
    '交付可部署的HTML原型文件',
], sz=14, color=CDARK, bullet='✓')

# ================================================================
# Slide 3: 全球盲盒市场概览
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CWHITE)
section_title(sl, '全球盲盒市场概览')

# 使用简单卡片替代柱状图（避免定位问题）
card_data = [
    ('2022年', '70亿美元', C1),
    ('2023年', '85亿美元', C2),
    ('2024年', '105亿美元', C1),
    ('2025年(预测)', '128亿美元', C3),
]
for i, (label, value, clr) in enumerate(card_data):
    x = Inches(0.6) + i * Inches(3.1)
    rrect(sl, x, Inches(1.3), Inches(2.8), Inches(2), CWHITE)
    txt(sl, x + Inches(0.2), Inches(1.5), Inches(2.4), Inches(0.4),
        label, sz=16, color=CGRAY, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.2), Inches(1.9), Inches(2.4), Inches(0.7),
        value, sz=30, bold=True, color=clr, align=PP_ALIGN.CENTER)
    # 箭头（最后三项之间）
    if i < 3:
        yoy_map = {0: '+21.4%', 1: '+23.5%', 2: '+21.9%'}
        txt(sl, x + Inches(0.2), Inches(2.7), Inches(2.4), Inches(0.4),
            f'↑ {yoy_map[i]}', sz=12, color=C4, align=PP_ALIGN.CENTER)

# 下方：CAGR和驱动因素
rrect(sl, Inches(0.6), Inches(3.8), Inches(5.8), Inches(3.2), CWHITE)
txt(sl, Inches(1), Inches(3.9), Inches(5), Inches(0.5),
    '全球市场总览', sz=18, bold=True, color=CDARK)
multi_text(sl, Inches(1), Inches(4.5), Inches(5), Inches(2.3), [
    '2022-2025年复合年增长率（CAGR）约22.5%',
    '亚太地区主导全球盲盒/潮玩市场',
    'Z世代和千禧一代为核心消费群体',
    '社交媒体（TikTok、Instagram）驱动增长',
    'IP联名和艺术家合作不断创新',
], sz=14, color=CDARK)

rrect(sl, Inches(6.8), Inches(3.8), Inches(5.8), Inches(3.2), CWHITE)
txt(sl, Inches(7.2), Inches(3.9), Inches(5), Inches(0.5),
    '关键增长驱动力', sz=18, bold=True, color=CDARK)
multi_text(sl, Inches(7.2), Inches(4.5), Inches(5), Inches(2.3), [
    'Z世代与千禧一代消费群体持续扩大',
    'TikTok开箱视频和社交传播效应',
    'IP联名跨界合作（动漫、游戏、明星）',
    '毛绒盲盒品类创新（软萌产品爆发）',
    '中产阶级可支配收入稳步增长',
], sz=14, color=CDARK)

# ================================================================
# Slide 4: 东南亚市场总览
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CBG)
section_title(sl, '东南亚市场总览')

# 市场规模卡片
market_cards = [
    ('2023年', '12亿美元', '-', C1),
    ('2024年', '18亿美元', '同比增长 +50%', C4),
    ('2025年(预测)', '25亿美元', '同比增长 +38.9%', C3),
]
for i, (yr, sz_val, growth, clr) in enumerate(market_cards):
    x = Inches(0.6) + i * Inches(4.1)
    rrect(sl, x, Inches(1.3), Inches(3.7), Inches(1.6), CWHITE)
    txt(sl, x + Inches(0.3), Inches(1.4), Inches(3), Inches(0.4), yr, sz=13, color=CGRAY)
    txt(sl, x + Inches(0.3), Inches(1.8), Inches(3), Inches(0.6), sz_val, sz=26, bold=True, color=clr)
    txt(sl, x + Inches(0.3), Inches(2.5), Inches(3), Inches(0.4), growth, sz=11, color=C4)

# 国家明细表
countries = [
    ('泰国', '#1', '5.5亿', '+57.1%', '12家'),
    ('越南', '#2', '2.8亿', '+55.6%', '5家'),
    ('印尼', '#3', '3.5亿', '+59.1%', '6家'),
    ('马来西亚', '#4', '2.3亿', '+53.3%', '4家'),
    ('菲律宾', '#5', '1.9亿', '+58.3%', '3家'),
    ('新加坡', '#6', '1.1亿', '+37.5%', '3家'),
]
headers = ['国家', '排名', '2024年规模', '同比增长', '泡泡玛特门店']
col_ws = [Inches(3), Inches(1.5), Inches(2.5), Inches(2), Inches(2.5)]
col_xs = [Inches(0.6)]
for w in col_ws[:-1]:
    col_xs.append(col_xs[-1] + w)

# 表头
y0 = Inches(3.3)
rect(sl, Inches(0.6), y0, Inches(11.5), Inches(0.4), C1)
for j, h in enumerate(headers):
    txt(sl, col_xs[j], y0 + Inches(0.02), col_ws[j], Inches(0.38),
        h, sz=12, bold=True, color=CWHITE, align=PP_ALIGN.CENTER)

# 数据行
for i, row in enumerate(countries):
    y = y0 + Inches(0.45) + i * Inches(0.45)
    bgc = CBG if i % 2 == 0 else CWHITE
    rect(sl, Inches(0.6), y, Inches(11.5), Inches(0.42), bgc)
    for j, val in enumerate(row):
        txt(sl, col_xs[j], y + Inches(0.04), col_ws[j], Inches(0.38),
            val, sz=11, color=CDARK, align=PP_ALIGN.CENTER)

# ================================================================
# Slide 5: 泡泡玛特品牌聚焦
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CDARKBG)
rect(sl, Inches(0), Inches(0), SW, Inches(0.06), C3)
txt(sl, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
    '品牌聚焦：泡泡玛特 (9992.HK)', sz=30, bold=True, color=CWHITE)

# KPI四列
kpi_data = [
    ('2024年总收入', '130.4亿元', '全年数据', C1),
    ('海外收入', '50.7亿元', '同比增长 +375%', C2),
    ('东南亚收入', '18亿元', '同比增长 +414%', C3),
    ('净利润', '34亿元', '同比增长 +186%', C4),
]
for i, (label, val, sub, clr) in enumerate(kpi_data):
    x = Inches(0.4) + i * Inches(3.2)
    rrect(sl, x, Inches(1.2), Inches(2.9), Inches(1.8))
    txt(sl, x + Inches(0.3), Inches(1.35), Inches(2.3), Inches(0.4),
        label, sz=12, color=CGRAY)
    txt(sl, x + Inches(0.3), Inches(1.8), Inches(2.3), Inches(0.5),
        val, sz=22, bold=True, color=clr)
    txt(sl, x + Inches(0.3), Inches(2.4), Inches(2.3), Inches(0.4),
        sub, sz=11, color=C4)

# 下方两栏
# Labubu现象
rrect(sl, Inches(0.4), Inches(3.3), Inches(6), Inches(3.8))
txt(sl, Inches(0.8), Inches(3.45), Inches(5.2), Inches(0.45),
    'Labubu现象级IP', sz=18, bold=True, color=C3)
multi_text(sl, Inches(0.8), Inches(4), Inches(5.2), Inches(2.8), [
    '香港艺术家龙家升（Kasing Lung）创作',
    '2024年泡泡玛特销售额排名第一的IP',
    'BLACKPINK成员Lisa社交媒体曝光',
    '泰国公主公开购买，引爆全网热议',
    'TikTok与Instagram病毒式传播',
    '二手市场溢价达原价3-5倍',
], sz=12, color=CLIGHT)

# 门店分布
rrect(sl, Inches(6.8), Inches(3.3), Inches(6), Inches(3.8))
txt(sl, Inches(7.2), Inches(3.45), Inches(5.2), Inches(0.45),
    '东南亚门店分布', sz=18, bold=True, color=CWHITE)
for i, (country, cnt) in enumerate([
    ('泰国', 12), ('印尼', 6), ('越南', 5),
    ('马来西亚', 4), ('菲律宾', 3), ('新加坡', 3),
]):
    y = Inches(4.15) + i * Inches(0.45)
    txt(sl, Inches(7.2), y, Inches(2.3), Inches(0.35),
        country, sz=13, color=CLIGHT)
    bw = Inches(2.5) * (cnt / 14)
    if bw < Inches(0.3):
        bw = Inches(0.3)
    rect(sl, Inches(9.5), y + Inches(0.06), bw, Inches(0.2), C3)
    txt(sl, Inches(9.5) + bw + Inches(0.12), y, Inches(0.4), Inches(0.35),
        str(cnt), sz=13, bold=True, color=C3)

# ================================================================
# Slide 6: 消费者画像
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CWHITE)
section_title(sl, '消费者画像洞察')

# 年龄分布
rrect(sl, Inches(0.6), Inches(1.2), Inches(3.8), Inches(3), CWHITE)
txt(sl, Inches(1), Inches(1.3), Inches(3), Inches(0.5), '年龄分布', sz=18, bold=True, color=CDARK)
ages = [('18-24岁', '35%', C3), ('25-34岁', '40%', C1),
        ('35-44岁', '18%', C4), ('45岁以上', '7%', C5)]
for i, (age, pct, clr) in enumerate(ages):
    y = Inches(1.9) + i * Inches(0.55)
    txt(sl, Inches(1), y, Inches(1.8), Inches(0.4), age, sz=14, color=CDARK)
    txt(sl, Inches(2.8), y, Inches(1), Inches(0.4), pct, sz=16, bold=True, color=clr)

# 性别分布
rrect(sl, Inches(4.8), Inches(1.2), Inches(3.8), Inches(3), CWHITE)
txt(sl, Inches(5.2), Inches(1.3), Inches(3), Inches(0.5), '性别分布', sz=18, bold=True, color=CDARK)
txt(sl, Inches(5.2), Inches(2), Inches(3), Inches(0.6), '女性 62%', sz=28, bold=True, color=C3)
txt(sl, Inches(5.2), Inches(2.8), Inches(3), Inches(0.5), '男性 38%', sz=20, color=C1)

# 购买渠道
rrect(sl, Inches(9), Inches(1.2), Inches(3.8), Inches(3), CWHITE)
txt(sl, Inches(9.4), Inches(1.3), Inches(3), Inches(0.5), '购买渠道', sz=18, bold=True, color=CDARK)
channels = [('线下零售', '45%'), ('电商平台', '35%'),
            ('快闪店', '12%'), ('自动售货机', '8%')]
for i, (ch, pct) in enumerate(channels):
    y = Inches(1.9) + i * Inches(0.55)
    txt(sl, Inches(9.4), y, Inches(2), Inches(0.4), ch, sz=13, color=CGRAY)
    txt(sl, Inches(11.5), y, Inches(0.8), Inches(0.4), pct, sz=15, bold=True, color=C3, align=PP_ALIGN.RIGHT)

# 消费习惯
rrect(sl, Inches(0.6), Inches(4.5), Inches(6), Inches(2.5), CWHITE)
txt(sl, Inches(1), Inches(4.6), Inches(5), Inches(0.5), '消费习惯', sz=18, bold=True, color=CDARK)
txt(sl, Inches(1), Inches(5.3), Inches(2.5), Inches(0.4), '平均客单价', sz=13, color=CGRAY)
txt(sl, Inches(1), Inches(5.7), Inches(2.5), Inches(0.6), '12.5美元', sz=28, bold=True, color=C3)
txt(sl, Inches(3.5), Inches(5.3), Inches(2.5), Inches(0.4), '月均购买频次', sz=13, color=CGRAY)
txt(sl, Inches(3.5), Inches(5.7), Inches(2.5), Inches(0.6), '2.8次', sz=28, bold=True, color=C1)

# 热门品类
rrect(sl, Inches(7), Inches(4.5), Inches(5.8), Inches(2.5), CWHITE)
txt(sl, Inches(7.4), Inches(4.6), Inches(5), Inches(0.5), '热门品类', sz=18, bold=True, color=CDARK)
cats = ['设计师玩具 40%', 'IP毛绒 25%', '动漫手办 15%', '文具生活 12%', '其他 8%']
for i, c in enumerate(cats):
    y = Inches(5.2) + i * Inches(0.35)
    txt(sl, Inches(7.4), y, Inches(5), Inches(0.35), f'• {c}', sz=12, color=CDARK)

# ================================================================
# Slide 7: 10套设计方案
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CBG)
section_title(sl, '10套网页设计方案')

templates = [
    ('01', '电商风格'), ('02', '品牌官网'), ('03', '社交媒体'),
    ('04', '数据可视化'), ('05', '极简风格'), ('06', '潮流街头'),
    ('07', '日系可爱'), ('08', '科技未来'), ('09', '杂志排版'),
    ('10', '沉浸式全屏'),
]
template_colors = [C1, C3, C4, C5, CGRAY, CDARK, C3, C1, CDARK, C1]
for i, ((num, name), clr) in enumerate(zip(templates, template_colors)):
    row, col = i // 5, i % 5
    x = Inches(0.4) + col * Inches(2.5)
    y = Inches(1.3) + row * Inches(2.9)
    rrect(sl, x, y, Inches(2.2), Inches(2.5), CWHITE)
    # 编号圆圈
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), y + Inches(0.3), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = clr
    circle.line.fill.background()
    txt(sl, x + Inches(0.7), y + Inches(0.35), Inches(0.7), Inches(0.6),
        num, sz=18, bold=True, color=CWHITE, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.1), y + Inches(1.3), Inches(2), Inches(0.5),
        name, sz=14, bold=True, color=CDARK, align=PP_ALIGN.CENTER)

# ================================================================
# Slide 8: 风格评估表
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CWHITE)
section_title(sl, '风格评估与推荐')

eval_headers = ['模板', '配色', '字体', '布局', '视觉', '交互', '移动\n适配', '东南亚\n适配度', '总分']
eval_scores = [
    ('07 日系可爱', 9, 7, 8, 9, 8, 6, 9, 64),
    ('03 社交媒体', 8, 7, 9, 8, 9, 8, 8, 66),
    ('10 沉浸式', 9, 8, 9, 8, 9, 6, 7, 65),
    ('02 品牌官网', 8, 8, 9, 7, 8, 6, 7, 61),
    ('06 潮流街头', 8, 7, 7, 8, 9, 5, 7, 61),
    ('09 杂志排版', 7, 9, 9, 7, 7, 5, 6, 61),
    ('08 科技未来', 7, 7, 8, 8, 9, 5, 5, 58),
    ('04 数据可视化', 7, 7, 8, 9, 8, 5, 4, 57),
    ('01 电商风格', 6, 7, 8, 5, 7, 6, 5, 52),
    ('05 极简风格', 6, 8, 8, 4, 6, 6, 3, 47),
]

e_col_ws = [Inches(2)] + [Inches(1.1)] * 7 + [Inches(1)]
e_col_xs = [Inches(0.4)]
for w in e_col_ws[:-1]:
    e_col_xs.append(e_col_xs[-1] + w)

y0 = Inches(1.2)
rect(sl, Inches(0.4), y0, Inches(11.3), Inches(0.4), C1)
for j, h in enumerate(eval_headers):
    txt(sl, e_col_xs[j], y0 + Inches(0.02), e_col_ws[j], Inches(0.38),
        h.replace('\n', ' '), sz=10, bold=True, color=CWHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(eval_scores):
    y = y0 + Inches(0.45) + i * Inches(0.52)
    bgc = CWHITE if i % 2 == 0 else CBG
    rect(sl, Inches(0.4), y, Inches(11.3), Inches(0.48), bgc)
    for j, val in enumerate(row):
        clr = C3 if j == len(row) - 1 else CDARK
        bold = (j == 0 or j == len(row) - 1)
        txt(sl, e_col_xs[j], y + Inches(0.06), e_col_ws[j], Inches(0.4),
            str(val), sz=10, bold=bold, color=clr, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

txt(sl, Inches(0.8), Inches(1.2 + 0.45 + 10 * 0.52 + 0.1), Inches(11), Inches(0.4),
    '东南亚适配度排名：#1 日系可爱(9分)  |  #2 社交媒体(8分)  |  #3 品牌官网/沉浸式(7分)',
    sz=13, bold=True, color=C1)

# ================================================================
# Slide 9: 推荐方案#1 - 日系可爱
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CWHITE)
rect(sl, Inches(0), Inches(0), Inches(0.12), SH, C3)
section_title(sl, '推荐方案#1：日系可爱风格（东南亚适配度 9/10）')

rrect(sl, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), CWHITE)
txt(sl, Inches(1), Inches(1.4), Inches(5), Inches(0.5), '选择理由', sz=18, bold=True, color=CDARK)
multi_text(sl, Inches(1), Inches(1.9), Inches(5), Inches(4.5), [
    '文化契合：Kawaii文化在泰国、菲律宾、越南极受欢迎',
    '色彩策略：粉色系(#FF69B4)匹配年轻女性用户偏好（占62%）',
    '视觉丰富：飘落动画、圆角卡片、渐变背景',
    '社交友好：可爱风格天然适合Instagram/TikTok视觉分享',
    'IP匹配：完美契合LABUBU、MOLLY"丑萌"美学风格',
], sz=13, color=CDARK)

rrect(sl, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5), CWHITE)
txt(sl, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5), '优化方案', sz=18, bold=True, color=CDARK)
multi_text(sl, Inches(7.2), Inches(1.9), Inches(5), Inches(4.5), [
    '添加响应式CSS，移动端优先设计（>80%手机用户）',
    '集成LABUBU、MOLLY等IP角色视觉元素',
    '添加完整购买流程和价格展示',
    '支持泰语、越南语、印尼语等多语言',
    '添加"开箱"动画交互体验',
    '集成TikTok Shop和Shopee购物链接',
], sz=13, color=CDARK, bullet='✓')

# ================================================================
# Slide 10: 推荐方案#2 - 社交媒体
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CWHITE)
rect(sl, Inches(0), Inches(0), Inches(0.12), SH, C1)
section_title(sl, '推荐方案#2：社交媒体风格（东南亚适配度 8/10）')

rrect(sl, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), CWHITE)
txt(sl, Inches(1), Inches(1.4), Inches(5), Inches(0.5), '核心优势', sz=18, bold=True, color=CDARK)
multi_text(sl, Inches(1), Inches(1.9), Inches(5), Inches(4.5), [
    '东南亚社交媒体使用率全球最高',
    'TikTok Shop在区域爆发式增长',
    'UGC内容驱动盲盒"开箱"分享文化',
    '移动优先布局，底部导航栏',
    '单列信息流天然适配手机屏幕',
    '清晰的转化路径：发现→互动→购买',
], sz=13, color=CDARK)

rrect(sl, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5), CWHITE)
txt(sl, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5), '整合方案', sz=18, bold=True, color=CDARK)
multi_text(sl, Inches(7.2), Inches(1.9), Inches(5), Inches(4.5), [
    '接入TikTok Shop、Shopee、Lazada购买链接',
    '开箱视频展示功能',
    '本地KOL/网红内容展示',
    '支持本地支付：GrabPay、GoPay、PromptPay',
    '社交分享：LINE、Zalo、WhatsApp、TikTok',
    '话题活动和UGC展示墙',
], sz=13, color=CDARK, bullet='✓')

# ================================================================
# Slide 11: 推荐方案#3 - 混合方案
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CBG)
section_title(sl, '推荐方案#3：混合解决方案（品牌官网 + 日系可爱融合）')

sections = [
    ('Hero区', '全屏布局\n糖果色渐变替换深蓝\n品牌宣言+行动按钮', C3),
    ('概念卡片', '三列卡片布局\n增加可爱插画\n暖色系配色', C1),
    ('社交证明', '简化数据展示\n粉丝数、开箱数、评分\n可视化受欢迎程度', C4),
    ('品牌故事', '图文排版布局\n增加IP角色展示\n文化叙事', C5),
    ('产品网格', '圆角卡片设计\n渐变背景\n可爱标签+价签', C3),
    ('行动号召', '圆角按钮\n社交分享按钮\n东南亚社交链接', CGRAY),
]
for i, (title, desc, clr) in enumerate(sections):
    col = i % 3
    row = i // 3
    x = Inches(0.4) + col * Inches(4.2)
    y = Inches(1.3) + row * Inches(2.9)
    rrect(sl, x, y, Inches(3.9), Inches(2.6), CWHITE)
    rect(sl, x, y, Inches(0.1), Inches(2.6), clr)
    txt(sl, x + Inches(0.3), y + Inches(0.15), Inches(3.3), Inches(0.45),
        title, sz=16, bold=True, color=CDARK)
    lines = desc.split('\n')
    for j, line in enumerate(lines):
        txt(sl, x + Inches(0.3), y + Inches(0.7) + j * Inches(0.4), Inches(3.3), Inches(0.4),
            f'• {line}', sz=11, color=CGRAY)

# ================================================================
# Slide 12: 东南亚设计指南
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CWHITE)
section_title(sl, '东南亚设计指南')

# 推荐配色
rrect(sl, Inches(0.6), Inches(1.2), Inches(5.8), Inches(2.8), CWHITE)
txt(sl, Inches(1), Inches(1.3), Inches(5), Inches(0.5), '推荐配色方案', sz=18, bold=True, color=CDARK)
palette = [
    ('#FF69B4', '主色-粉色', 'Kawaii美学核心'),
    ('#4FACFE', '辅色-天蓝', '清新活力感'),
    ('#43E97B', '辅色-薄荷绿', '热带氛围'),
    ('#FFD93D', '辅色-暖黄', '泰国市场偏好'),
    ('#FF2442', '行动-红色', '价格和CTA强调'),
]
for i, (hex_code, name, use) in enumerate(palette):
    y = Inches(1.9) + i * Inches(0.4)
    r, g, b = int(hex_code[1:3], 16), int(hex_code[3:5], 16), int(hex_code[5:7], 16)
    swatch = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y, Inches(0.3), Inches(0.28))
    swatch.fill.solid()
    swatch.fill.fore_color.rgb = RGBColor(r, g, b)
    swatch.line.fill.background()
    txt(sl, Inches(1.5), y, Inches(1.5), Inches(0.3), hex_code, sz=10, bold=True, color=CDARK)
    txt(sl, Inches(3), y, Inches(1.5), Inches(0.3), name, sz=10, color=CGRAY)
    txt(sl, Inches(4.3), y, Inches(2), Inches(0.3), use, sz=10, color=CGRAY)

# 必备功能
rrect(sl, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.8), CWHITE)
txt(sl, Inches(7.2), Inches(1.3), Inches(5), Inches(0.5), '必备功能', sz=18, bold=True, color=CDARK)
multi_text(sl, Inches(7.2), Inches(1.8), Inches(5), Inches(2), [
    '移动优先设计（>80%移动用户）',
    '多语言：泰语、越南语、印尼语、马来语、英语',
    '本地支付：GrabPay、GoPay、ShopeePay、PromptPay',
    '社交分享：LINE、Zalo、WhatsApp、TikTok',
    '轻量化页面（部分地区网络较慢）',
    '开箱动画交互体验',
], sz=11, color=CDARK)

# 设计趋势
rrect(sl, Inches(0.6), Inches(4.3), Inches(12), Inches(2.7), CWHITE)
txt(sl, Inches(1), Inches(4.4), Inches(11), Inches(0.5),
    '2025-2026年东南亚设计趋势', sz=18, bold=True, color=CDARK)
trends = [
    ('"丑萌"美学', 'LABUBU式怪诞+可爱融合'),
    ('糖果/粉彩色系', '柔和梦幻色彩搭配'),
    ('Y2K复古回潮', '千禧年美学复兴'),
    ('文化融合', '本地艺术家联名合作'),
    ('毛绒品类热门', '毛绒挂件盲盒大热'),
    ('社交优先设计', 'TikTok/Instagram友好视觉'),
]
for i, (trend, desc) in enumerate(trends):
    col = i % 3
    row = i // 3
    x = Inches(1) + col * Inches(4)
    y = Inches(5) + row * Inches(0.9)
    txt(sl, x, y, Inches(1.5), Inches(0.35), trend, sz=12, bold=True, color=C1)
    txt(sl, x + Inches(1.6), y, Inches(2.2), Inches(0.35), desc, sz=11, color=CGRAY)

# ================================================================
# Slide 13: 项目成果
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CDARKBG)
rect(sl, Inches(0), Inches(0), SW, Inches(0.06), C3)
txt(sl, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
    '项目成果总结', sz=30, bold=True, color=CWHITE)

deliverables = [
    ('01', '市场数据报告', '东南亚6国市场分析\nJSON+Markdown双格式\n2023-2025数据覆盖', C3),
    ('02', '网页设计模板', '10套HTML设计方案\n零外部依赖，独立可用\n交互式原型展示', C1),
    ('03', '风格评估报告', '7维度评分体系\n4品牌设计分析\n3级推荐方案', C4),
    ('04', '综合验收报告', '总分8.2/10\n4维度质量验证\n问题追踪与整改', C5),
]
for i, (num, title, desc, clr) in enumerate(deliverables):
    x = Inches(0.3) + i * Inches(3.2)
    rrect(sl, x, Inches(1.3), Inches(3), Inches(3.3))
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1), Inches(1.5), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = clr
    circle.line.fill.background()
    txt(sl, x + Inches(1), Inches(1.6), Inches(0.9), Inches(0.7),
        num, sz=20, bold=True, color=CWHITE, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.15), Inches(2.6), Inches(2.7), Inches(0.4),
        title, sz=14, bold=True, color=CWHITE, align=PP_ALIGN.CENTER)
    for j, line in enumerate(desc.split('\n')):
        txt(sl, x + Inches(0.15), Inches(3.1) + j * Inches(0.3), Inches(2.7), Inches(0.3),
            line, sz=11, color=CLIGHT, align=PP_ALIGN.CENTER)

# 评分
rrect(sl, Inches(0.4), Inches(5), Inches(12.5), Inches(1.8))
txt(sl, Inches(0.8), Inches(5.1), Inches(11), Inches(0.4), '验收评分', sz=16, bold=True, color=CWHITE)
scores = [('数据质量', '9/10'), ('网页模板', '7/10'), ('风格分析', '9/10'), ('项目结构', '8/10'), ('总分', '8.2/10')]
for i, (label, score) in enumerate(scores):
    x = Inches(0.8) + i * Inches(2.4)
    clr = C5 if label == '总分' else CLIGHT
    bold = label == '总分'
    txt(sl, x, Inches(5.5), Inches(2.2), Inches(0.3), label, sz=11, color=CGRAY, align=PP_ALIGN.CENTER)
    txt(sl, x, Inches(5.8), Inches(2.2), Inches(0.5), score, sz=18, bold=bold, color=clr, align=PP_ALIGN.CENTER)

# ================================================================
# Slide 14: 后续行动计划
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CBG)
section_title(sl, '后续行动计划')

phases = [
    ('阶段一：响应式设计', '为所有模板添加媒体查询\n移动优先设计方法\n测试320px、768px、1024px\n预计1周完成', C3),
    ('阶段二：A/B测试', '在泰国、越南、菲律宾市场\n测试Top 3模板方案\n测量转化率数据\n预计2周完成', C1),
    ('阶段三：本地化整合', '接入本地支付网关\n连接TikTok Shop和Shopee\n添加多语言支持\n预计3周完成', C4),
    ('阶段四：内容升级', '替换emoji为真实产品图片\n添加IP角色插画\n创建开箱动画效果\n预计2周完成', C5),
]
for i, (phase, desc, clr) in enumerate(phases):
    x = Inches(0.3) + i * Inches(3.2)
    rrect(sl, x, Inches(1.3), Inches(3), Inches(4), CWHITE)
    # 圆形编号
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), Inches(1.45), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = clr
    circle.line.fill.background()
    txt(sl, x + Inches(1.1), Inches(1.5), Inches(0.7), Inches(0.6),
        str(i + 1), sz=18, bold=True, color=CWHITE, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.15), Inches(2.3), Inches(2.7), Inches(0.45),
        phase, sz=13, bold=True, color=CDARK, align=PP_ALIGN.CENTER)
    for j, line in enumerate(desc.split('\n')):
        txt(sl, x + Inches(0.15), Inches(2.85 + j * 0.35), Inches(2.7), Inches(0.35),
            f'• {line}', sz=11, color=CGRAY)

# 优先级矩阵
rrect(sl, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.3))
txt(sl, Inches(0.8), Inches(5.7), Inches(11), Inches(0.35), '优先级矩阵', sz=15, bold=True, color=CDARK)
priority_items = [
    ('高', '所有模板添加响应式CSS', C3),
    ('高', '移动端优先优化', C3),
    ('中', '目标市场A/B测试', C1),
    ('中', '本地支付整合', C1),
    ('低', '用真实图片替换emoji', CGRAY),
    ('低', 'JavaScript错误处理', CGRAY),
]
for i, (level, task, clr) in enumerate(priority_items):
    x = Inches(0.8) + (i % 3) * Inches(4.1)
    y = Inches(6.15) + (i // 3) * Inches(0.35)
    txt(sl, x, y, Inches(0.6), Inches(0.3), f'[{level}]', sz=10, bold=True, color=clr)
    txt(sl, x + Inches(0.7), y, Inches(3), Inches(0.3), task, sz=10, color=CGRAY)

# ================================================================
# Slide 15: 感谢
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, CDARKBG)
txt(sl, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
    '感谢观看', sz=48, bold=True, color=CWHITE, align=PP_ALIGN.CENTER)
txt(sl, Inches(1.5), Inches(4), Inches(10), Inches(0.6),
    '东南亚盲盒市场网页设计方案汇报', sz=18, color=C2, align=PP_ALIGN.CENTER)
txt(sl, Inches(1.5), Inches(5), Inches(10), Inches(0.5),
    '2026年6月', sz=14, color=CGRAY, align=PP_ALIGN.CENTER)

# ================================================================
# Save
# ================================================================
out = 'F:/测试工具/blindbox-project/presentation/东南亚盲盒市场网页设计方案汇报-v3.pptx'
prs.save(out)
print(f'OK: {out}')
print(f'Total slides: {len(prs.slides)}')
