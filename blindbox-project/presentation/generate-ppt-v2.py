"""
生成东南亚盲盒市场网页设计方案汇报PPT（中文专业版）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 专业配色方案
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)      # 主色-深蓝
PRIMARY_LIGHT = RGBColor(0x4F, 0x86, 0xF7) # 主色-浅蓝
ACCENT = RGBColor(0xFF, 0x6B, 0x6B)       # 强调色-珊瑚红
ACCENT2 = RGBColor(0x4E, 0xC9, 0xB0)      # 强调色-青绿
ACCENT3 = RGBColor(0xFF, 0xB8, 0x4C)      # 强调色-金黄
DARK = RGBColor(0x2C, 0x3E, 0x50)         # 深色文字
GRAY = RGBColor(0x7F, 0x8C, 0x8D)         # 灰色文字
LIGHT_GRAY = RGBColor(0xBD, 0xC3, 0xC7)   # 浅灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)     # 浅背景
DARK_BG = RGBColor(0x1E, 0x2D, 0x3D)      # 深背景

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    w = width or prs.slide_width
    h = height or prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT,
                 font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, bold=False, color=DARK,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6),
                  font_name='Microsoft YaHei'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_card(slide, left, top, width, height, fill_color=WHITE,
             border_color=None, shadow=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if shadow:
        shape.shadow.inherit = False
    return shape


def add_section_header(slide, title, subtitle=None):
    """添加统一的章节标题"""
    # 顶部装饰条
    add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), prs.slide_width, Inches(0.08))

    # 标题
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                 title, font_size=32, bold=True, color=PRIMARY)

    # 副标题（仅中文）
    if subtitle and not any(c.isascii() and c.isalpha() for c in subtitle):
        add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=GRAY)


# ============================================================
# 幻灯片1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

# 装饰元素
for cx, cy, sz, clr in [
    (Inches(1), Inches(1), Inches(2), PRIMARY_LIGHT),
    (Inches(11), Inches(5.5), Inches(1.5), ACCENT),
    (Inches(10), Inches(0.5), Inches(1), ACCENT2),
    (Inches(2), Inches(6), Inches(0.8), ACCENT3),
]:
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, sz, sz)
    c.fill.solid()
    c.fill.fore_color.rgb = clr
    c.line.fill.background()

# 主标题
add_text_box(slide, Inches(1.5), Inches(2), Inches(10), Inches(1.5),
             '东南亚盲盒市场网页设计方案', font_size=44, bold=True,
             color=WHITE, alignment=PP_ALIGN.CENTER)

# 副标题
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1),
             '市场分析 | 设计方案 | 风格评估', font_size=24, color=PRIMARY_LIGHT,
             alignment=PP_ALIGN.CENTER)

# 日期和版本
add_text_box(slide, Inches(1.5), Inches(5), Inches(10), Inches(0.5),
             '2026年6月', font_size=16, color=GRAY,
             alignment=PP_ALIGN.CENTER)

# 公司信息
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             '市场调研与设计团队', font_size=14, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)


# ============================================================
# 幻灯片2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, '汇报目录')

# 目录项
toc_items = [
    ('01', '项目背景与目标', '市场机会与项目定位'),
    ('02', '全球盲盒市场概览', '市场规模与增长趋势'),
    ('03', '东南亚市场深度分析', '六国市场对比与机会'),
    ('04', '品牌格局分析', '泡泡玛特、名创优品等'),
    ('05', '消费者画像洞察', '年龄、性别、消费习惯'),
    ('06', '网页设计方案展示', '10套风格模板详解'),
    ('07', '风格评估与推荐', '7维度评分与推荐方案'),
    ('08', '项目成果总结', '交付物与验收评分'),
    ('09', '后续行动计划', '分阶段实施建议'),
]

for i, (num, title, desc) in enumerate(toc_items):
    row = i // 3
    col = i % 3
    x = Inches(0.8) + col * Inches(4)
    y = Inches(1.5) + row * Inches(1.8)

    # 编号圆圈
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY if i % 2 == 0 else ACCENT
    circle.line.fill.background()
    add_text_box(slide, x, y + Inches(0.05), Inches(0.6), Inches(0.5),
                 num, font_size=18, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # 标题
    add_text_box(slide, x + Inches(0.7), y, Inches(3), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK)

    # 描述
    add_text_box(slide, x + Inches(0.7), y + Inches(0.4), Inches(3), Inches(0.3),
                 desc, font_size=12, color=GRAY)


# ============================================================
# 幻灯片3: 项目背景与目标
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section_header(slide, '项目背景与目标')

# 左侧-市场机会
card = add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2))
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
             '市场机会', font_size=20, bold=True, color=PRIMARY)
txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(5), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
items = [
    '东南亚盲盒市场年复合增长率：28.5%',
    '2025年市场规模预计达25亿美元',
    '泡泡玛特海外收入突破50亿元',
    '泰国、印尼、越南为核心增长引擎',
    '年轻人口结构（中位年龄30岁）',
    '社交媒体渗透率全球领先',
]
for i, item in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  • {item}'
    p.font.size = Pt(14)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(8)

# 右侧-项目目标
card = add_card(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(5.2))
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5), Inches(0.5),
             '项目目标', font_size=20, bold=True, color=ACCENT)
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(2.2), Inches(5), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
items = [
    '深度分析东南亚盲盒市场数据',
    '设计10套网页模板概念方案',
    '评估品牌风格适配性',
    '推荐最佳市场解决方案',
    '提供可落地的设计指南',
    '交付可部署的HTML原型',
]
for i, item in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  ✓ {item}'
    p.font.size = Pt(14)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(8)


# ============================================================
# 幻灯片4: 全球盲盒市场概览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, '全球盲盒市场概览')

# 柱状图模拟
years = ['2022年', '2023年', '2024年', '2025年(预测)']
values = [70, 85, 105, 128]
colors = [LIGHT_GRAY, PRIMARY_LIGHT, PRIMARY, ACCENT]
bar_left = Inches(1.5)
bar_bottom = Inches(5.5)
bar_width = Inches(1.8)
max_val = 140

for i, (year, val, clr) in enumerate(zip(years, values, colors)):
    x = bar_left + i * Inches(2.8)
    h_ratio = val / max_val
    bar_h = Inches(3.5) * h_ratio
    bar_y = bar_bottom - bar_h
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, bar_y, bar_width, bar_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = clr
    shape.line.fill.background()
    # 数值标签
    add_text_box(slide, x, bar_y - Inches(0.4), bar_width, Inches(0.4),
                 f'{val}亿美元', font_size=18, bold=True, color=clr,
                 alignment=PP_ALIGN.CENTER)
    # 年份标签
    add_text_box(slide, x, bar_bottom + Inches(0.1), bar_width, Inches(0.4),
                 year, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# CAGR标签
add_text_box(slide, Inches(9.5), Inches(2), Inches(3), Inches(1.5),
             '年复合增长率\n22.5%', font_size=28, bold=True, color=PRIMARY,
             alignment=PP_ALIGN.CENTER)

# 关键驱动因素
add_text_box(slide, Inches(9.5), Inches(4), Inches(3), Inches(0.4),
             '核心驱动因素', font_size=16, bold=True, color=DARK)
drivers = ['Z世代消费崛起', '社交媒体传播', 'IP联名创新', '毛绒品类爆发']
txBox = slide.shapes.add_textbox(Inches(9.5), Inches(4.5), Inches(3), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
for i, d in enumerate(drivers):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  • {d}'
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft YaHei'


# ============================================================
# 幻灯片5: 东南亚市场深度分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section_header(slide, '东南亚市场深度分析')

# 市场规模卡片
data = [
    ('2023年', '12亿美元', '-', PRIMARY),
    ('2024年', '18亿美元', '+50%', ACCENT2),
    ('2025年(预测)', '25亿美元', '+38.9%', ACCENT),
]
for i, (year, size, growth, clr) in enumerate(data):
    x = Inches(0.8) + i * Inches(4)
    card = add_card(slide, x, Inches(1.5), Inches(3.5), Inches(1.8))
    add_text_box(slide, x + Inches(0.3), Inches(1.6), Inches(3), Inches(0.5),
                 year, font_size=14, color=GRAY)
    add_text_box(slide, x + Inches(0.3), Inches(2.1), Inches(3), Inches(0.6),
                 size, font_size=28, bold=True, color=clr)
    if growth != '-':
        add_text_box(slide, x + Inches(0.3), Inches(2.8), Inches(3), Inches(0.4),
                     f'同比增长: {growth}', font_size=12, color=ACCENT2)

# 国家对比表（按市场规模排名）
countries = [
    ('泰国', '#1', '5.5亿美元', '+57.1%', '12家门店'),
    ('越南', '#2', '2.8亿美元', '+55.6%', '5家门店'),
    ('印尼', '#3', '3.5亿美元', '+59.1%', '6家门店'),
    ('马来西亚', '#4', '2.3亿美元', '+53.3%', '4家门店'),
    ('菲律宾', '#5', '1.9亿美元', '+58.3%', '3家门店'),
    ('新加坡', '#6', '1.1亿美元', '+37.5%', '3家门店'),
]
headers = ['国家', '排名', '2024年规模', '同比增长', '泡泡玛特门店']
y_start = Inches(3.8)
col_widths = [Inches(2.5), Inches(1.5), Inches(2), Inches(2), Inches(2.5)]
col_x = [Inches(0.8)]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

# 表头
hdr_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), y_start, Inches(10.5), Inches(0.4))
hdr_shape.fill.solid()
hdr_shape.fill.fore_color.rgb = PRIMARY
hdr_shape.line.fill.background()
for j, h in enumerate(headers):
    add_text_box(slide, col_x[j], y_start, col_widths[j], Inches(0.4),
                 h, font_size=12, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

for i, row in enumerate(countries):
    y = y_start + Inches(0.5) + i * Inches(0.4)
    bg_clr = LIGHT_BG if i % 2 == 0 else WHITE
    row_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(10.5), Inches(0.4))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = bg_clr
    row_bg.line.fill.background()
    for j, val in enumerate(row):
        add_text_box(slide, col_x[j], y, col_widths[j], Inches(0.4),
                     val, font_size=11, color=DARK,
                     alignment=PP_ALIGN.CENTER)


# ============================================================
# 幻灯片6: 泡泡玛特品牌聚焦
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_section_header(slide, '品牌聚焦：泡泡玛特 (9992.HK)')

# KPI卡片
kpis = [
    ('总收入', '130.4亿元', '2024全年'),
    ('海外收入', '50.7亿元', '+375% 同比'),
    ('东南亚收入', '18亿元', '+414% 同比'),
    ('净利润', '34亿元', '+186% 同比'),
]
for i, (label, value, sub) in enumerate(kpis):
    x = Inches(0.5) + i * Inches(3.1)
    card = add_card(slide, x, Inches(1.5), Inches(2.8), Inches(2))
    add_text_box(slide, x + Inches(0.2), Inches(1.6), Inches(2.4), Inches(0.4),
                 label, font_size=12, color=GRAY)
    add_text_box(slide, x + Inches(0.2), Inches(2.1), Inches(2.4), Inches(0.6),
                 value, font_size=20, bold=True, color=ACCENT)
    add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(2.4), Inches(0.4),
                 sub, font_size=11, color=ACCENT2)

# Labubu现象
card = add_card(slide, Inches(0.5), Inches(4), Inches(5.8), Inches(3))
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(5), Inches(0.5),
             'Labubu现象级IP', font_size=18, bold=True, color=ACCENT)
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(5.2), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
labubu_points = [
    '香港艺术家Kasing Lung创作',
    '2024年泡泡玛特销售额第一IP',
    'BLACKPINK Lisa社交媒体曝光',
    '泰国公主公开购买',
    'TikTok/Instagram病毒式传播',
    '二手市场溢价3-5倍',
]
for i, pt in enumerate(labubu_points):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  • {pt}'
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(6)

# 门店分布
card = add_card(slide, Inches(6.8), Inches(4), Inches(5.8), Inches(3))
add_text_box(slide, Inches(7.1), Inches(4.1), Inches(5), Inches(0.5),
             '东南亚门店分布', font_size=18, bold=True, color=WHITE)
stores = [
    ('泰国', '12'), ('印尼', '6'), ('越南', '5'),
    ('马来西亚', '4'), ('菲律宾', '3'), ('新加坡', '3'),
]
for i, (country, count) in enumerate(stores):
    y = Inches(4.8) + i * Inches(0.35)
    add_text_box(slide, Inches(7.1), y, Inches(2.5), Inches(0.35),
                 country, font_size=12, color=LIGHT_GRAY)
    bar_w = Inches(2.5) * (int(count) / 14)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), y + Inches(0.05), bar_w, Inches(0.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text_box(slide, Inches(9.5) + bar_w + Inches(0.1), y, Inches(0.5), Inches(0.35),
                 count, font_size=11, bold=True, color=ACCENT)


# ============================================================
# 幻灯片7: 消费者画像洞察
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, '消费者画像洞察')

# 年龄分布
card = add_card(slide, Inches(0.5), Inches(1.3), Inches(4), Inches(5.5))
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(3.5), Inches(0.5),
             '年龄分布', font_size=16, bold=True, color=DARK)
ages = [('18-24岁', '35%', ACCENT), ('25-34岁', '40%', PRIMARY),
        ('35-44岁', '18%', ACCENT2), ('45岁以上', '7%', ACCENT3)]
for i, (age, pct, clr) in enumerate(ages):
    y = Inches(2.1) + i * Inches(1)
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1), y + Inches(0.05), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = clr
    dot.line.fill.background()
    add_text_box(slide, Inches(1.5), y, Inches(1.5), Inches(0.4),
                 age, font_size=14, color=DARK)
    add_text_box(slide, Inches(3), y, Inches(1), Inches(0.4),
                 pct, font_size=16, bold=True, color=clr)

# 性别分布
card = add_card(slide, Inches(4.8), Inches(1.3), Inches(3.8), Inches(2.5))
add_text_box(slide, Inches(5.1), Inches(1.4), Inches(3.3), Inches(0.5),
             '性别分布', font_size=16, bold=True, color=DARK)
add_text_box(slide, Inches(5.1), Inches(2), Inches(3.3), Inches(0.6),
             '女性 62%', font_size=24, bold=True, color=ACCENT)
add_text_box(slide, Inches(5.1), Inches(2.7), Inches(3.3), Inches(0.5),
             '男性 38%', font_size=18, color=PRIMARY)

# 购买渠道
card = add_card(slide, Inches(4.8), Inches(4.1), Inches(3.8), Inches(2.7))
add_text_box(slide, Inches(5.1), Inches(4.2), Inches(3.3), Inches(0.5),
             '购买渠道', font_size=16, bold=True, color=DARK)
channels = [('线下零售', '45%'), ('电商平台', '35%'),
            ('快闪店', '12%'), ('自动售货机', '8%')]
for i, (ch, pct) in enumerate(channels):
    y = Inches(4.8) + i * Inches(0.45)
    add_text_box(slide, Inches(5.1), y, Inches(2), Inches(0.4),
                 ch, font_size=12, color=GRAY)
    add_text_box(slide, Inches(7.2), y, Inches(1), Inches(0.4),
                 pct, font_size=14, bold=True, color=ACCENT)

# 消费习惯
card = add_card(slide, Inches(9), Inches(1.3), Inches(3.8), Inches(5.5))
add_text_box(slide, Inches(9.3), Inches(1.4), Inches(3.3), Inches(0.5),
             '消费习惯', font_size=16, bold=True, color=DARK)
add_text_box(slide, Inches(9.3), Inches(2.2), Inches(3.3), Inches(0.4),
             '平均客单价', font_size=12, color=GRAY)
add_text_box(slide, Inches(9.3), Inches(2.6), Inches(3.3), Inches(0.6),
             '12.5美元', font_size=28, bold=True, color=ACCENT)
add_text_box(slide, Inches(9.3), Inches(3.3), Inches(3.3), Inches(0.4),
             '月均购买频次', font_size=12, color=GRAY)
add_text_box(slide, Inches(9.3), Inches(3.7), Inches(3.3), Inches(0.6),
             '2.8次', font_size=28, bold=True, color=PRIMARY)
add_text_box(slide, Inches(9.3), Inches(4.5), Inches(3.3), Inches(0.4),
             '热门品类', font_size=12, bold=True, color=DARK)
cats = ['设计师玩具 40%', 'IP毛绒 25%', '动漫手办 15%',
        '文具生活 12%', '其他 8%']
txBox = slide.shapes.add_textbox(Inches(9.3), Inches(5), Inches(3.3), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, c in enumerate(cats):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  • {c}'
    p.font.size = Pt(11)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft YaHei'


# ============================================================
# 幻灯片8: 10套网页设计方案
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section_header(slide, '10套网页设计方案', '10 WEB DESIGN TEMPLATES')

templates = [
    ('01', '电商风格', PRIMARY), ('02', '品牌官网', ACCENT),
    ('03', '社交媒体', ACCENT2), ('04', '数据可视化', ACCENT3),
    ('05', '极简风格', GRAY), ('06', '潮流街头', DARK),
    ('07', '日系可爱', ACCENT), ('08', '科技未来', PRIMARY),
    ('09', '杂志排版', DARK), ('10', '沉浸式全屏', PRIMARY),
]
for i, (num, name, clr) in enumerate(templates):
    row = i // 5
    col = i % 5
    x = Inches(0.5) + col * Inches(2.5)
    y = Inches(1.2) + row * Inches(3)
    card = add_card(slide, x, y, Inches(2.2), Inches(2.5), border_color=clr)
    # 编号徽章
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.7), y + Inches(0.3), Inches(0.8), Inches(0.8))
    badge.fill.solid()
    badge.fill.fore_color.rgb = clr
    badge.line.fill.background()
    add_text_box(slide, x + Inches(0.7), y + Inches(0.4), Inches(0.8), Inches(0.6),
                 num, font_size=18, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.3), Inches(2), Inches(0.5),
                 name, font_size=14, bold=True, color=DARK,
                 alignment=PP_ALIGN.CENTER)


# ============================================================
# 幻灯片9: 风格评估与推荐
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, '风格评估与推荐')

# 评分表
headers = ['模板', '配色', '字体', '布局', '视觉', '交互', '移动适配', '东南亚适配', '总分']
scores = [
    ('07 日系可爱', 9, 7, 8, 9, 8, 6, 9, 64),
    ('03 社交媒体', 8, 7, 9, 8, 9, 8, 8, 66),
    ('10 沉浸式', 9, 8, 9, 8, 9, 6, 7, 65),
    ('02 品牌官网', 8, 8, 9, 7, 8, 6, 7, 61),
    ('06 潮流街头', 8, 7, 7, 8, 9, 5, 7, 61),
    ('09 杂志排版', 7, 9, 9, 7, 7, 5, 6, 61),
    ('08 科技未来', 7, 7, 8, 8, 9, 5, 5, 58),
    ('04 数据可视化', 7, 7, 8, 9, 8, 5, 4, 57),
    ('01 电商风格', 6, 7, 8, 5, 7, 6, 5, 52),
    ('05 极简风格', 6, 8, 8, 4, 6, 6, 3, 47),
]

y_start = Inches(1.2)
col_w = [Inches(1.8)] + [Inches(1.1)] * 8
col_x = [Inches(0.4)]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

# 表头
hdr_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.4), y_start, Inches(11.5), Inches(0.45))
hdr_bg.fill.solid()
hdr_bg.fill.fore_color.rgb = PRIMARY
hdr_bg.line.fill.background()
for j, h in enumerate(headers):
    add_text_box(slide, col_x[j], y_start + Inches(0.02), col_w[j], Inches(0.4),
                 h, font_size=10, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

# 数据行
for i, row in enumerate(scores):
    y = y_start + Inches(0.5) + i * Inches(0.55)
    bg_clr = WHITE if i % 2 == 0 else LIGHT_BG
    row_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), y, Inches(11.5), Inches(0.5))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = bg_clr
    row_bg.line.fill.background()
    for j, val in enumerate(row):
        clr = ACCENT if j == 8 else DARK
        bld = j == 0 or j == 8
        sz = 12 if j == 0 else 11
        add_text_box(slide, col_x[j], y + Inches(0.05), col_w[j], Inches(0.4),
                     str(val), font_size=sz, bold=bld, color=clr,
                     alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# 推荐说明
add_text_box(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
             '东南亚适配度Top 3: #1 日系可爱(9分) | #2 社交媒体(8分) | #3 品牌官网/沉浸式(7分)',
             font_size=14, bold=True, color=PRIMARY)


# ============================================================
# 幻灯片10: 推荐方案1-日系可爱风格
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, ACCENT, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
add_section_header(slide, '推荐方案#1：日系可爱风格')

# 优势分析
card = add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2))
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.5),
             '为什么选择这个风格', font_size=18, bold=True, color=DARK)
reasons = [
    '文化契合：Kawaii文化在泰国、菲律宾、越南极受欢迎',
    '色彩策略：粉色系(#FF69B4)匹配年轻女性用户偏好',
    '视觉丰富：飘落爱心/星星动画、圆角卡片、渐变背景',
    '社交友好：可爱风格天然适合Instagram/TikTok分享',
    'IP匹配：完美契合LABUBU、MOLLY"丑萌"美学',
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.2), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
for i, r in enumerate(reasons):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  • {r}'
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(8)

# 优化建议
card = add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5),
             '优化方案', font_size=18, bold=True, color=DARK)
opts = [
    '添加响应式CSS，移动端优先设计',
    '集成LABUBU、MOLLY IP角色视觉',
    '添加购买流程和价格展示',
    '支持泰语、越南语、印尼语',
    '添加"开箱"动画交互',
    '集成TikTok Shop / Shopee链接',
]
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
for i, o in enumerate(opts):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  ✓ {o}'
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(8)


# ============================================================
# 幻灯片11: 推荐方案2-社交媒体风格
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
add_section_header(slide, '推荐方案#2：社交媒体风格')

# 左侧-优势
card = add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2))
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.5),
             '核心优势', font_size=18, bold=True, color=DARK)
strengths = [
    '东南亚社交媒体使用率全球最高',
    'TikTok Shop在区域爆发式增长',
    'UGC内容驱动盲盒"开箱"文化',
    '移动优先布局，底部导航栏',
    '单列信息流天然适配手机屏幕',
    '清晰转化路径：发现→互动→购买',
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.2), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
for i, s in enumerate(strengths):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  • {s}'
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(8)

# 右侧-整合方案
card = add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5),
             '整合方案', font_size=18, bold=True, color=DARK)
integrations = [
    'TikTok Shop / Shopee / Lazada购买链接',
    '开箱视频展示功能',
    '本地KOL/网红内容展示',
    '本地支付：GrabPay、GoPay、PromptPay',
    '社交分享：LINE、Zalo、WhatsApp、TikTok',
    '话题活动和UGC展示墙',
]
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(integrations):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  ✓ {item}'
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(8)


# ============================================================
# 幻灯片12: 推荐方案3-混合方案
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section_header(slide, '推荐方案#3：混合解决方案')

sections = [
    ('Hero区', '来自模板02的全屏布局\n糖果色渐变替换深蓝\n品牌宣言+行动按钮', ACCENT),
    ('概念卡片', '模板02的三列卡片布局\n增加可爱插画+emoji\n暖色系配色', PRIMARY),
    ('社交证明', '简化数据展示\n粉丝数、开箱数、评分\n可视化受欢迎程度', ACCENT2),
    ('品牌故事', '模板02的图文排版\n增加IP角色展示\n文化叙事', ACCENT3),
    ('产品网格', '模板07的圆角卡片\n渐变背景\n可爱标签+价签', ACCENT),
    ('行动号召', '模板07的圆角按钮\n社交分享按钮\n东南亚社交链接', GRAY),
]
for i, (title, desc, clr) in enumerate(sections):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(1.6) + row * Inches(2.8)
    card = add_card(slide, x, y, Inches(3.9), Inches(2.5))
    # 色彩条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, Inches(0.12), Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.1), Inches(3.4), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK)
    txBox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.6), Inches(3.4), Inches(1.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f'  {line}'
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.font.name = 'Microsoft YaHei'


# ============================================================
# 幻灯片13: 设计指南
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, '东南亚设计指南')

# 配色方案
card = add_card(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(2.8))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.4),
             '推荐配色方案', font_size=16, bold=True, color=DARK)
colors_info = [
    ('#FF69B4', '主色-粉色', 'Kawaii美学核心'),
    ('#4FACFE', '辅助色-天蓝', '清新活力'),
    ('#43E97B', '辅助色-薄荷绿', '热带氛围'),
    ('#FFD93D', '强调色-暖黄', '泰国市场偏好'),
    ('#FF2442', '行动色-红色', '价格/CTA强调'),
]
for i, (hex_val, name, use) in enumerate(colors_info):
    y = Inches(1.9) + i * Inches(0.4)
    clr_parts = hex_val.lstrip('#')
    r, g, b = int(clr_parts[0:2], 16), int(clr_parts[2:4], 16), int(clr_parts[4:6], 16)
    swatch = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y, Inches(0.35), Inches(0.3))
    swatch.fill.solid()
    swatch.fill.fore_color.rgb = RGBColor(r, g, b)
    swatch.line.fill.background()
    add_text_box(slide, Inches(1.5), y, Inches(1.5), Inches(0.3),
                 f'{hex_val}', font_size=10, bold=True, color=DARK)
    add_text_box(slide, Inches(3), y, Inches(1.5), Inches(0.3),
                 name, font_size=10, color=GRAY)
    add_text_box(slide, Inches(4.3), y, Inches(2), Inches(0.3),
                 use, font_size=10, color=GRAY)

# 必备功能
card = add_card(slide, Inches(7), Inches(1.2), Inches(5.8), Inches(2.8))
add_text_box(slide, Inches(7.3), Inches(1.3), Inches(5), Inches(0.4),
             '必备功能', font_size=16, bold=True, color=DARK)
features = [
    '移动优先设计（>80%移动用户）',
    '多语言支持：泰语、越南语、印尼语、马来语、英语',
    '本地支付：GrabPay、GoPay、ShopeePay、PromptPay',
    '社交分享：LINE、Zalo、WhatsApp、TikTok',
    '轻量化页面（部分地区网络较慢）',
    '开箱动画交互',
]
txBox = slide.shapes.add_textbox(Inches(7.3), Inches(1.9), Inches(5.3), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
for i, f in enumerate(features):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  • {f}'
    p.font.size = Pt(11)
    p.font.color.rgb = DARK
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(4)

# 设计趋势
card = add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8))
add_text_box(slide, Inches(0.8), Inches(4.4), Inches(11), Inches(0.4),
             '2025-2026东南亚设计趋势', font_size=16, bold=True, color=DARK)

trends = [
    ('"丑萌"美学', 'LABUBU式怪诞+可爱融合'),
    ('糖果/粉彩色系', '柔和梦幻色彩搭配'),
    ('Y2K复古回潮', '千禧年美学复兴'),
    ('文化融合', '本地艺术家联名合作'),
    ('毛绒/软质材料', '毛绒挂件盲盒大热'),
    ('社交优先设计', 'TikTok/Instagram友好视觉'),
]
for i, (trend, desc) in enumerate(trends):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(5) + row * Inches(1)
    add_text_box(slide, x, y, Inches(3.8), Inches(0.4),
                 trend, font_size=13, bold=True, color=PRIMARY)
    add_text_box(slide, x, y + Inches(0.35), Inches(3.8), Inches(0.4),
                 desc, font_size=11, color=GRAY)


# ============================================================
# 幻灯片14: 项目成果总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_section_header(slide, '项目成果总结')

deliverables = [
    ('市场数据报告', '东南亚6国市场分析\nJSON+Markdown双格式\n2023-2025数据覆盖', ACCENT),
    ('网页设计模板', '10套HTML设计方案\n零外部依赖\n交互式原型', PRIMARY),
    ('风格评估报告', '7维度评分体系\n4品牌分析\n3级推荐方案', ACCENT2),
    ('综合验收报告', '总分8.2/10\n4维度验证\n质量问题追踪', ACCENT3),
]
for i, (title, desc, clr) in enumerate(deliverables):
    x = Inches(0.3) + i * Inches(3.2)
    card = add_card(slide, x, Inches(1.3), Inches(3), Inches(3.5))
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(1), Inches(1.5), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = clr
    circle.line.fill.background()
    num = ['01', '02', '03', '04'][i]
    add_text_box(slide, x + Inches(1), Inches(1.6), Inches(1), Inches(0.8),
                 num, font_size=20, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.6), Inches(0.5),
                 title, font_size=14, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    txBox = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.2), Inches(2.6), Inches(1.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f'  {line}'
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER

# 验收评分
card = add_card(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8))
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(0.4),
             '验收评分', font_size=16, bold=True, color=WHITE)
score_items = [
    ('数据质量', '9/10', '30%'), ('网页模板', '7/10', '35%'),
    ('风格分析', '9/10', '25%'), ('项目结构', '8/10', '10%'),
    ('总分', '8.2/10', '100%'),
]
for i, (label, score, weight) in enumerate(score_items):
    x = Inches(0.8) + i * Inches(2.4)
    clr = ACCENT3 if label == '总分' else LIGHT_GRAY
    bld = label == '总分'
    add_text_box(slide, x, Inches(5.8), Inches(2.2), Inches(0.3),
                 label, font_size=11, color=GRAY,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.1), Inches(2.2), Inches(0.5),
                 score, font_size=18, bold=bld, color=clr,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.5), Inches(2.2), Inches(0.3),
                 f'权重: {weight}', font_size=9, color=GRAY,
                 alignment=PP_ALIGN.CENTER)


# ============================================================
# 幻灯片15: 后续行动计划
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section_header(slide, '后续行动计划')

steps = [
    ('阶段1：响应式设计', '为所有模板添加媒体查询\n移动优先设计方法\n测试320px、768px、1024px断点', ACCENT),
    ('阶段2：A/B测试', '在目标市场测试Top 3模板\n测量转化率数据\n重点测试泰国、越南、菲律宾', PRIMARY),
    ('阶段3：本地化整合', '接入本地支付网关\n连接TikTok Shop / Shopee\n添加多语言支持', ACCENT2),
    ('阶段4：内容与素材', '用真实产品图替换emoji\n添加IP角色插画\n创建开箱动画效果', ACCENT3),
]
for i, (phase, desc, clr) in enumerate(steps):
    x = Inches(0.5) + i * Inches(3.1)
    card = add_card(slide, x, Inches(1.2), Inches(2.9), Inches(4))
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(1), Inches(1.4), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = clr
    circle.line.fill.background()
    add_text_box(slide, x + Inches(1), Inches(1.5), Inches(0.8), Inches(0.6),
                 str(i + 1), font_size=20, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.4), Inches(2.5), Inches(0.5),
                 phase, font_size=14, bold=True, color=DARK,
                 alignment=PP_ALIGN.CENTER)
    txBox = slide.shapes.add_textbox(x + Inches(0.2), Inches(3), Inches(2.5), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f'  {line}'
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.font.name = 'Microsoft YaHei'
        p.space_before = Pt(6)

# 优先级矩阵
card = add_card(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5))
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.4),
             '优先级矩阵', font_size=14, bold=True, color=DARK)
priorities = [
    ('高优先级', '所有模板添加响应式CSS', ACCENT),
    ('高优先级', '移动端优先优化', ACCENT),
    ('中优先级', '目标市场A/B测试', PRIMARY),
    ('中优先级', '本地支付整合', PRIMARY),
    ('低优先级', '用真实图片替换emoji', GRAY),
    ('低优先级', 'JavaScript错误处理', GRAY),
]
for i, (level, task, clr) in enumerate(priorities):
    x = Inches(0.8) + (i % 3) * Inches(4.1)
    y = Inches(6.1) + (i // 3) * Inches(0.4)
    add_text_box(slide, x, y, Inches(0.8), Inches(0.35),
                 f'[{level}]', font_size=10, bold=True, color=clr)
    add_text_box(slide, x + Inches(0.8), y, Inches(3), Inches(0.35),
                 task, font_size=10, color=GRAY)


# ============================================================
# 幻灯片16: 感谢页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

# 装饰元素
for cx, cy, sz, clr in [
    (Inches(2), Inches(2), Inches(1.5), ACCENT),
    (Inches(10), Inches(5), Inches(2), PRIMARY_LIGHT),
    (Inches(11), Inches(1), Inches(1), ACCENT2),
    (Inches(1), Inches(6), Inches(0.8), ACCENT3),
]:
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, sz, sz)
    c.fill.solid()
    c.fill.fore_color.rgb = clr
    c.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
             '感谢观看', font_size=48, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4), Inches(10), Inches(0.6),
             '东南亚盲盒市场网页设计方案汇报', font_size=18, color=PRIMARY_LIGHT,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5), Inches(10), Inches(0.5),
             '2026年6月', font_size=14, color=GRAY,
             alignment=PP_ALIGN.CENTER)

# 联系信息
add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
             '市场调研与设计团队', font_size=12, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)


# ============================================================
# 保存文件
# ============================================================
output_path = 'F:/测试工具/blindbox-project/presentation/东南亚盲盒市场网页设计方案汇报.pptx'
prs.save(output_path)
print(f'PPT已生成: {output_path}')
print(f'共{len(prs.slides)}页幻灯片')
