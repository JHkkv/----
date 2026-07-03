"""
Generate PPT presentation for Southeast Asia Blind Box Market project.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette
PINK = RGBColor(0xFF, 0x69, 0xB4)
DARK_PINK = RGBColor(0xFF, 0x24, 0x42)
BLUE = RGBColor(0x4F, 0xAC, 0xFE)
GREEN = RGBColor(0x43, 0xE9, 0x7B)
YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_PINK_BG = RGBColor(0xFF, 0xF5, 0xF8)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    w = width or prs.slide_width
    h = height or prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, bold=False, color=BLACK,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6),
                  font_name='Microsoft YaHei'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_card(slide, left, top, width, height, fill_color=WHITE,
             border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ============================================================
# Slide 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# Decorative circles
for cx, cy, sz, clr in [
    (Inches(1), Inches(1), Inches(2), PINK),
    (Inches(11), Inches(5.5), Inches(1.5), BLUE),
    (Inches(10), Inches(0.5), Inches(1), GREEN),
    (Inches(2), Inches(6), Inches(0.8), YELLOW),
]:
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, sz, sz)
    c.fill.solid()
    c.fill.fore_color.rgb = clr
    c.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(2), Inches(10), Inches(1.5),
             'Southeast Asia Blind Box Market', font_size=44, bold=True,
             color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1),
             'Web Design Proposal', font_size=32, color=PINK,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.8),
             'Market Analysis  |  Design Concepts  |  Style Evaluation',
             font_size=18, color=RGBColor(0xAA, 0xAA, 0xAA),
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6), Inches(10), Inches(0.5),
             '2026.06', font_size=16, color=RGBColor(0x88, 0x88, 0x88),
             alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: Project Background
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_PINK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Project Background & Objectives', font_size=36, bold=True,
             color=DARK_PINK, alignment=PP_ALIGN.LEFT)

# Left column
card = add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2))
txBox = add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5),
                     'Market Opportunity', font_size=22, bold=True, color=BLACK)
txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(5), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
items = [
    'SE Asia blind box market CAGR: 28.5% (2023-2028)',
    'Market size: $12B (2023) -> $25B (2025E)',
    'Pop Mart overseas revenue exceeded RMB 5B',
    'Thailand, Indonesia, Vietnam as key drivers',
    'Young demographics (median age 30)',
    'Social media penetration driving purchases',
]
for i, item in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  {item}'
    p.font.size = Pt(15)
    p.font.color.rgb = GRAY
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(8)

# Right column
card = add_card(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(5.2))
txBox = add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5),
                     'Project Objectives', font_size=22, bold=True, color=BLACK)
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
items = [
    'Analyze SE Asia blind box market data',
    'Design 10 web template concepts',
    'Evaluate brand style alignment',
    'Recommend best-fit solutions',
    'Provide actionable design guidelines',
    'Deliver deployable HTML prototypes',
]
for i, item in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  {item}'
    p.font.size = Pt(15)
    p.font.color.rgb = GRAY
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(8)

# ============================================================
# Slide 3: Global Market Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Global Blind Box Market Overview', font_size=36, bold=True,
             color=DARK_PINK)

# Bar chart simulation with rectangles
years = ['2022', '2023', '2024', '2025E']
values = [70, 85, 105, 128]
colors = [BLUE, GREEN, PINK, DARK_PINK]
bar_left = Inches(1.5)
bar_bottom = Inches(5.5)
bar_width = Inches(1.8)
max_val = 140

for i, (year, val, clr) in enumerate(zip(years, values, colors)):
    x = bar_left + i * Inches(2.8)
    h_ratio = val / max_val
    bar_h = Inches(3.5) * h_ratio
    bar_y = bar_bottom - bar_h
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, bar_y, bar_width, bar_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = clr
    shape.line.fill.background()
    # Value label
    add_text_box(slide, x, bar_y - Inches(0.4), bar_width, Inches(0.4),
                 f'${val}B', font_size=20, bold=True, color=clr,
                 alignment=PP_ALIGN.CENTER)
    # Year label
    add_text_box(slide, x, bar_bottom + Inches(0.1), bar_width, Inches(0.4),
                 year, font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# CAGR label
add_text_box(slide, Inches(9.5), Inches(2), Inches(3), Inches(1.5),
             'CAGR\n22.5%', font_size=36, bold=True, color=DARK_PINK,
             alignment=PP_ALIGN.CENTER)

# Key drivers
add_text_box(slide, Inches(9.5), Inches(4), Inches(3), Inches(0.4),
             'Key Drivers', font_size=16, bold=True, color=BLACK)
drivers = ['Gen Z & Millennials', 'Social Media (TikTok)',
           'IP Collaborations', 'Plush Toy Innovation']
txBox = slide.shapes.add_textbox(Inches(9.5), Inches(4.5), Inches(3), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
for i, d in enumerate(drivers):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  {d}'
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    p.font.name = 'Microsoft YaHei'

# ============================================================
# Slide 4: SE Asia Market Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_PINK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Southeast Asia Market Overview', font_size=36, bold=True,
             color=DARK_PINK)

# Market size cards
data = [
    ('2023', '$1.2B', '-', BLUE),
    ('2024', '$1.8B', '+50%', GREEN),
    ('2025E', '$2.5B', '+38.9%', PINK),
]
for i, (year, size, growth, clr) in enumerate(data):
    x = Inches(0.8) + i * Inches(4)
    card = add_card(slide, x, Inches(1.5), Inches(3.5), Inches(1.8))
    add_text_box(slide, x + Inches(0.3), Inches(1.6), Inches(3), Inches(0.5),
                 year, font_size=18, color=GRAY)
    add_text_box(slide, x + Inches(0.3), Inches(2.1), Inches(3), Inches(0.6),
                 size, font_size=36, bold=True, color=clr)
    if growth != '-':
        add_text_box(slide, x + Inches(0.3), Inches(2.8), Inches(3), Inches(0.4),
                     f'YoY: {growth}', font_size=14, color=GREEN)

# Country breakdown table
countries = [
    ('Thailand', '#1', '$550M', '+57.1%', '12 stores'),
    ('Indonesia', '#2', '$350M', '+59.1%', '6 stores'),
    ('Vietnam', '#3', '$280M', '+55.6%', '5 stores'),
    ('Malaysia', '#4', '$230M', '+53.3%', '4 stores'),
    ('Philippines', '#5', '$190M', '+58.3%', '3 stores'),
    ('Singapore', '#6', '$110M', '+37.5%', '3 stores'),
]
headers = ['Country', 'Rank', '2024 Size', 'YoY Growth', 'Pop Mart Stores']
y_start = Inches(3.8)
col_widths = [Inches(2.5), Inches(1.5), Inches(2), Inches(2), Inches(2.5)]
col_x = [Inches(0.8)]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

# Header row
for j, h in enumerate(headers):
    add_text_box(slide, col_x[j], y_start, col_widths[j], Inches(0.4),
                 h, font_size=14, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
# Header bg
hdr_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), y_start, Inches(10.5), Inches(0.4))
hdr_shape.fill.solid()
hdr_shape.fill.fore_color.rgb = DARK_PINK
hdr_shape.line.fill.background()

for i, row in enumerate(countries):
    y = y_start + Inches(0.5) + i * Inches(0.4)
    bg_clr = LIGHT_PINK_BG if i % 2 == 0 else WHITE
    row_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(10.5), Inches(0.4))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = bg_clr
    row_bg.line.fill.background()
    for j, val in enumerate(row):
        add_text_box(slide, col_x[j], y, col_widths[j], Inches(0.4),
                     val, font_size=13, color=BLACK,
                     alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 5: Country Details
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Country Market Details', font_size=36, bold=True, color=DARK_PINK)

country_data = [
    ('Thailand', '$550M', '+57.1%', '12', 'Labubu, MOLLY, SKULLPANDA',
     'Pop Mart #1 overseas market\nLisa (BLACKPINK) drove demand\n3-5x resale premium'),
    ('Indonesia', '$350M', '+59.1%', '6', 'LABUBU, MOLLY, DIMOO',
     'Largest population (270M+)\nTikTok Shop key channel\nRising middle class'),
    ('Vietnam', '$280M', '+55.6%', '5', 'MOLLY, DIMOO, LABUBU',
     'Young demographics\nHigh social media penetration\nShopee/Lazada active'),
]

for i, (name, size, growth, stores, ips, notes) in enumerate(country_data):
    x = Inches(0.5) + i * Inches(4.2)
    card = add_card(slide, x, Inches(1.3), Inches(3.9), Inches(5.5))

    # Country header
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(3.9), Inches(0.7))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = [PINK, BLUE, GREEN][i]
    hdr.line.fill.background()
    add_text_box(slide, x + Inches(0.2), Inches(1.35), Inches(3.5), Inches(0.6),
                 name, font_size=22, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Metrics
    add_text_box(slide, x + Inches(0.3), Inches(2.2), Inches(3.3), Inches(0.5),
                 f'Market Size: {size}  |  Growth: {growth}', font_size=14,
                 bold=True, color=BLACK)
    add_text_box(slide, x + Inches(0.3), Inches(2.7), Inches(3.3), Inches(0.4),
                 f'Pop Mart Stores: {stores}', font_size=13, color=GRAY)
    add_text_box(slide, x + Inches(0.3), Inches(3.1), Inches(3.3), Inches(0.4),
                 f'Hot IPs: {ips}', font_size=13, color=PINK)

    # Notes
    txBox = slide.shapes.add_textbox(x + Inches(0.3), Inches(3.7), Inches(3.3), Inches(2.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(notes.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f'  {line}'
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.font.name = 'Microsoft YaHei'
        p.space_before = Pt(6)

# ============================================================
# Slide 6: Pop Mart Brand Spotlight
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Brand Spotlight: Pop Mart (9992.HK)', font_size=36, bold=True,
             color=WHITE)

# KPI cards
kpis = [
    ('Total Revenue', 'RMB 13.04B', '2024 Full Year'),
    ('Overseas Revenue', 'RMB 5.07B', '+375% YoY'),
    ('SE Asia Revenue', 'RMB 1.8B', '+414% YoY'),
    ('Net Profit', 'RMB 3.4B', '+186% YoY'),
]
for i, (label, value, sub) in enumerate(kpis):
    x = Inches(0.5) + i * Inches(3.1)
    card = add_card(slide, x, Inches(1.5), Inches(2.8), Inches(2))
    add_text_box(slide, x + Inches(0.2), Inches(1.6), Inches(2.4), Inches(0.4),
                 label, font_size=14, color=GRAY)
    add_text_box(slide, x + Inches(0.2), Inches(2.1), Inches(2.4), Inches(0.6),
                 value, font_size=24, bold=True, color=PINK)
    add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(2.4), Inches(0.4),
                 sub, font_size=13, color=GREEN)

# Labubu phenomenon
card = add_card(slide, Inches(0.5), Inches(4), Inches(5.8), Inches(3))
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(5), Inches(0.5),
             'Labubu Phenomenon', font_size=22, bold=True, color=PINK)
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(5.2), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
labubu_points = [
    'Created by artist Kasing Lung (Hong Kong)',
    '#1 selling IP for Pop Mart in 2024',
    'BLACKPINK Lisa spotted with Labubu products',
    'Thai princess endorsement',
    'Social media viral spread on TikTok/Instagram',
    'Resale market premium: 3-5x original price',
]
for i, pt in enumerate(labubu_points):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  {pt}'
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(6)

# Store distribution
card = add_card(slide, Inches(6.8), Inches(4), Inches(5.8), Inches(3))
add_text_box(slide, Inches(7.1), Inches(4.1), Inches(5), Inches(0.5),
             'SE Asia Store Distribution', font_size=22, bold=True, color=WHITE)
stores = [
    ('Thailand', '12'), ('Indonesia', '6'), ('Vietnam', '5'),
    ('Malaysia', '4'), ('Philippines', '3'), ('Singapore', '3'),
]
for i, (country, count) in enumerate(stores):
    y = Inches(4.8) + i * Inches(0.35)
    add_text_box(slide, Inches(7.1), y, Inches(2.5), Inches(0.35),
                 country, font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC))
    # Bar
    bar_w = Inches(2.5) * (int(count) / 14)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), y + Inches(0.05), bar_w, Inches(0.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PINK
    bar.line.fill.background()
    add_text_box(slide, Inches(9.5) + bar_w + Inches(0.1), y, Inches(0.5), Inches(0.35),
                 count, font_size=13, bold=True, color=PINK)

# ============================================================
# Slide 7: Consumer Profile
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_PINK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Consumer Profile', font_size=36, bold=True, color=DARK_PINK)

# Age distribution - donut simulation
card = add_card(slide, Inches(0.5), Inches(1.3), Inches(4), Inches(5.5))
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(3.5), Inches(0.5),
             'Age Distribution', font_size=20, bold=True, color=BLACK)
ages = [('18-24', '35%', PINK), ('25-34', '40%', DARK_PINK),
        ('35-44', '18%', BLUE), ('45+', '7%', GREEN)]
for i, (age, pct, clr) in enumerate(ages):
    y = Inches(2.1) + i * Inches(1)
    # Color indicator
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1), y + Inches(0.05), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = clr
    dot.line.fill.background()
    add_text_box(slide, Inches(1.5), y, Inches(1.5), Inches(0.4),
                 age, font_size=16, color=BLACK)
    add_text_box(slide, Inches(3), y, Inches(1), Inches(0.4),
                 pct, font_size=18, bold=True, color=clr)

# Gender
card = add_card(slide, Inches(4.8), Inches(1.3), Inches(3.8), Inches(2.5))
add_text_box(slide, Inches(5.1), Inches(1.4), Inches(3.3), Inches(0.5),
             'Gender Split', font_size=20, bold=True, color=BLACK)
add_text_box(slide, Inches(5.1), Inches(2), Inches(3.3), Inches(0.6),
             'Female 62%', font_size=28, bold=True, color=PINK)
add_text_box(slide, Inches(5.1), Inches(2.7), Inches(3.3), Inches(0.5),
             'Male 38%', font_size=22, color=BLUE)

# Purchase channels
card = add_card(slide, Inches(4.8), Inches(4.1), Inches(3.8), Inches(2.7))
add_text_box(slide, Inches(5.1), Inches(4.2), Inches(3.3), Inches(0.5),
             'Purchase Channels', font_size=20, bold=True, color=BLACK)
channels = [('Offline Retail', '45%'), ('E-commerce', '35%'),
            ('Pop-up Stores', '12%'), ('Vending', '8%')]
for i, (ch, pct) in enumerate(channels):
    y = Inches(4.8) + i * Inches(0.45)
    add_text_box(slide, Inches(5.1), y, Inches(2), Inches(0.4),
                 ch, font_size=14, color=GRAY)
    add_text_box(slide, Inches(7.2), y, Inches(1), Inches(0.4),
                 pct, font_size=16, bold=True, color=PINK)

# Spending habits
card = add_card(slide, Inches(9), Inches(1.3), Inches(3.8), Inches(5.5))
add_text_box(slide, Inches(9.3), Inches(1.4), Inches(3.3), Inches(0.5),
             'Spending Habits', font_size=20, bold=True, color=BLACK)
add_text_box(slide, Inches(9.3), Inches(2.2), Inches(3.3), Inches(0.4),
             'Avg. Spend', font_size=14, color=GRAY)
add_text_box(slide, Inches(9.3), Inches(2.6), Inches(3.3), Inches(0.6),
             '$12.50', font_size=32, bold=True, color=DARK_PINK)
add_text_box(slide, Inches(9.3), Inches(3.3), Inches(3.3), Inches(0.4),
             'Monthly Frequency', font_size=14, color=GRAY)
add_text_box(slide, Inches(9.3), Inches(3.7), Inches(3.3), Inches(0.6),
             '2.8x', font_size=32, bold=True, color=BLUE)
add_text_box(slide, Inches(9.3), Inches(4.5), Inches(3.3), Inches(0.4),
             'Top Categories', font_size=14, bold=True, color=BLACK)
cats = ['Designer Toy 40%', 'IP Plush 25%', 'Anime Figure 15%',
        'Stationery 12%', 'Other 8%']
txBox = slide.shapes.add_textbox(Inches(9.3), Inches(5), Inches(3.3), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, c in enumerate(cats):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  {c}'
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    p.font.name = 'Microsoft YaHei'

# ============================================================
# Slide 8: 10 Template Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '10 Web Design Templates', font_size=36, bold=True, color=DARK_PINK)

templates = [
    ('01', 'E-commerce', BLUE), ('02', 'Brand Official', PINK),
    ('03', 'Social Media', DARK_PINK), ('04', 'Data Viz', GREEN),
    ('05', 'Minimalist', GRAY), ('06', 'Street Style', BLACK),
    ('07', 'Kawaii Cute', PINK), ('08', 'Sci-Fi Future', BLUE),
    ('09', 'Magazine', GRAY), ('10', 'Immersive', DARK_PINK),
]
for i, (num, name, clr) in enumerate(templates):
    row = i // 5
    col = i % 5
    x = Inches(0.5) + col * Inches(2.5)
    y = Inches(1.2) + row * Inches(3)
    card = add_card(slide, x, y, Inches(2.2), Inches(2.5), border_color=clr)
    # Number badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.7), y + Inches(0.3), Inches(0.8), Inches(0.8))
    badge.fill.solid()
    badge.fill.fore_color.rgb = clr
    badge.line.fill.background()
    add_text_box(slide, x + Inches(0.7), y + Inches(0.4), Inches(0.8), Inches(0.6),
                 num, font_size=22, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.3), Inches(2), Inches(0.5),
                 name, font_size=16, bold=True, color=BLACK,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.8), Inches(2), Inches(0.5),
                 f'template-{num}', font_size=11, color=GRAY,
                 alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 9: Template Scoring
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_PINK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             'Template Evaluation Scores', font_size=36, bold=True, color=DARK_PINK)

# Scoring table
headers = ['Template', 'Color', 'Font', 'Layout', 'Visual', 'Interact', 'Mobile', 'SEA Fit', 'Total']
scores = [
    ('07 Kawaii', 9, 7, 8, 9, 8, 6, 9, 64),
    ('03 Social', 8, 7, 9, 8, 9, 8, 8, 66),
    ('10 Immersive', 9, 8, 9, 8, 9, 6, 7, 65),
    ('02 Brand', 8, 8, 9, 7, 8, 6, 7, 61),
    ('06 Street', 8, 7, 7, 8, 9, 5, 7, 61),
    ('09 Magazine', 7, 9, 9, 7, 7, 5, 6, 61),
    ('08 Sci-Fi', 7, 7, 8, 8, 9, 5, 5, 58),
    ('04 Data Viz', 7, 7, 8, 9, 8, 5, 4, 57),
    ('01 E-com', 6, 7, 8, 5, 7, 6, 5, 52),
    ('05 Minimal', 6, 8, 8, 4, 6, 6, 3, 47),
]

y_start = Inches(1.2)
col_w = [Inches(1.8)] + [Inches(1.1)] * 8
col_x = [Inches(0.4)]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

# Header
hdr_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.4), y_start, Inches(11.5), Inches(0.45))
hdr_bg.fill.solid()
hdr_bg.fill.fore_color.rgb = DARK_PINK
hdr_bg.line.fill.background()
for j, h in enumerate(headers):
    add_text_box(slide, col_x[j], y_start + Inches(0.02), col_w[j], Inches(0.4),
                 h, font_size=12, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

# Rows
for i, row in enumerate(scores):
    y = y_start + Inches(0.5) + i * Inches(0.55)
    bg_clr = WHITE if i % 2 == 0 else LIGHT_PINK_BG
    row_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), y, Inches(11.5), Inches(0.5))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = bg_clr
    row_bg.line.fill.background()
    for j, val in enumerate(row):
        clr = DARK_PINK if j == 8 else BLACK
        bld = j == 0 or j == 8
        sz = 14 if j == 0 else 13
        add_text_box(slide, col_x[j], y + Inches(0.05), col_w[j], Inches(0.4),
                     str(val), font_size=sz, bold=bld, color=clr,
                     alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# Highlight note
add_text_box(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
             'Top 3 by SEA Fit: #1 Kawaii (9)  |  #2 Social Media (8)  |  #3 Brand Official / Immersive (7)',
             font_size=15, bold=True, color=DARK_PINK)

# ============================================================
# Slide 10: Recommendation 1 - Kawaii
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
# Pink accent bar
add_shape_bg(slide, PINK, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(11), Inches(0.7),
             'Recommendation #1: Kawaii Style (Template 07)', font_size=32,
             bold=True, color=DARK_PINK)
add_text_box(slide, Inches(0.5), Inches(1), Inches(3), Inches(0.4),
             'SEA Fit Score: 9/10', font_size=20, bold=True, color=PINK)

# Why this works
card = add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2))
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.5),
             'Why This Works', font_size=22, bold=True, color=BLACK)
reasons = [
    'Cultural alignment: Kawaii culture extremely popular in Thailand, Philippines, Vietnam, Malaysia',
    'Color strategy: Pink (#FF69B4) + soft gradients match young female user preferences',
    'Visual richness: Falling hearts/stars animations, rounded cards, gradient backgrounds',
    'Social-friendly: Cute style naturally suits Instagram/TikTok visual sharing',
    'IP match: Perfectly fits LABUBU, MOLLY "ugly-cute" aesthetic',
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.2), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
for i, r in enumerate(reasons):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  {r}'
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(10)

# Optimization suggestions
card = add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5),
             'Optimization Plan', font_size=22, bold=True, color=BLACK)
opts = [
    'Add responsive CSS for mobile-first design',
    'Integrate LABUBU, MOLLY IP character visuals',
    'Add purchase flow and price display',
    'Support Thai, Vietnamese, Indonesian languages',
    'Add "box opening" animation interaction',
    'Integrate TikTok Shop / Shopee links',
]
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
for i, o in enumerate(opts):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  {o}'
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(10)

# ============================================================
# Slide 11: Recommendation 2 - Social Media
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, DARK_PINK, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(11), Inches(0.7),
             'Recommendation #2: Social Media Style (Template 03)', font_size=32,
             bold=True, color=DARK_PINK)
add_text_box(slide, Inches(0.5), Inches(1), Inches(3), Inches(0.4),
             'SEA Fit Score: 8/10', font_size=20, bold=True, color=DARK_PINK)

# Left column
card = add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2))
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.5),
             'Key Strengths', font_size=22, bold=True, color=BLACK)
strengths = [
    'SE Asia has highest social media usage globally',
    'TikTok Shop explosive growth in the region',
    'UGC content drives blind box "unboxing" culture',
    'Mobile-first layout with bottom nav bar',
    'Single-column feed naturally fits phone screens',
    'Clear conversion: discover -> engage -> purchase',
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.2), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
for i, s in enumerate(strengths):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  {s}'
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(10)

# Right column
card = add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2))
add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.5),
             'Integration Plan', font_size=22, bold=True, color=BLACK)
integrations = [
    'TikTok Shop / Shopee / Lazada purchase links',
    'Unboxing video showcase functionality',
    'Local KOL/influencer content display',
    'Local payment: GrabPay, GoPay, PromptPay',
    'Social sharing: LINE, Zalo, WhatsApp, TikTok',
    'Hashtag campaigns and UGC gallery',
]
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(integrations):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  {item}'
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(10)

# ============================================================
# Slide 12: Recommendation 3 - Hybrid
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_PINK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             'Recommendation #3: Hybrid Solution', font_size=36, bold=True,
             color=DARK_PINK)
add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(0.4),
             'Brand Official (Template 02) + Kawaii Style (Template 07) Fusion',
             font_size=18, color=GRAY)

# Fusion diagram
sections = [
    ('Hero Section', 'Full-screen layout from T02\nCandy-color gradient (replace dark blue)\nBrand statement + CTA', PINK),
    ('Concept Cards', '3-column card layout from T02\nAdd cute illustrations + emoji\nWarm color scheme', BLUE),
    ('Social Proof', 'Simplified data display\nFan count, unbox count, rating\nVisual proof of popularity', GREEN),
    ('Brand Story', 'Image-text layout from T02\nAdd IP character showcase\nCultural narrative', YELLOW),
    ('Product Grid', 'Rounded cards from T07\nGradient backgrounds\nCute labels + price tags', DARK_PINK),
    ('CTA + Footer', 'Rounded buttons from T07\nSocial share buttons\nSE Asia social links (LINE, Zalo)', GRAY),
]
for i, (title, desc, clr) in enumerate(sections):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(1.6) + row * Inches(2.8)
    card = add_card(slide, x, y, Inches(3.9), Inches(2.5))
    # Color bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, Inches(0.12), Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.1), Inches(3.4), Inches(0.4),
                 title, font_size=18, bold=True, color=BLACK)
    txBox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.6), Inches(3.4), Inches(1.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f'  {line}'
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.font.name = 'Microsoft YaHei'

# ============================================================
# Slide 13: Design Guidelines
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             'SE Asia Design Guidelines', font_size=36, bold=True, color=DARK_PINK)

# Color palette
card = add_card(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(2.8))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.4),
             'Color Palette', font_size=20, bold=True, color=BLACK)
colors_info = [
    ('#FF69B4', 'Primary Pink', 'Kawaii aesthetic'),
    ('#4FACFE', 'Sky Blue', 'Fresh & energetic'),
    ('#43E97B', 'Mint Green', 'Tropical vibe'),
    ('#FFD93D', 'Warm Yellow', 'Thai market preference'),
    ('#FF2442', 'Action Red', 'Price/CTA emphasis'),
]
for i, (hex_val, name, use) in enumerate(colors_info):
    y = Inches(1.9) + i * Inches(0.4)
    clr_parts = hex_val.lstrip('#')
    r, g, b = int(clr_parts[0:2], 16), int(clr_parts[2:4], 16), int(clr_parts[4:6], 16)
    swatch = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y, Inches(0.35), Inches(0.3))
    swatch.fill.solid()
    swatch.fill.fore_color.rgb = RGBColor(r, g, b)
    swatch.line.fill.background()
    add_text_box(slide, Inches(1.5), y, Inches(1.5), Inches(0.3),
                 f'{hex_val}', font_size=12, bold=True, color=BLACK)
    add_text_box(slide, Inches(3), y, Inches(1.5), Inches(0.3),
                 name, font_size=12, color=GRAY)
    add_text_box(slide, Inches(4.3), y, Inches(2), Inches(0.3),
                 use, font_size=12, color=GRAY)

# Must-have features
card = add_card(slide, Inches(7), Inches(1.2), Inches(5.8), Inches(2.8))
add_text_box(slide, Inches(7.3), Inches(1.3), Inches(5), Inches(0.4),
             'Must-Have Features', font_size=20, bold=True, color=BLACK)
features = [
    'Mobile-first design (>80% mobile users)',
    'Multi-language: Thai, Vietnamese, Indonesian, Malay, English',
    'Local payments: GrabPay, GoPay, ShopeePay, PromptPay',
    'Social sharing: LINE, Zalo, WhatsApp, TikTok',
    'Lightweight pages (slow network in some areas)',
    'Box-opening animation interaction',
]
txBox = slide.shapes.add_textbox(Inches(7.3), Inches(1.9), Inches(5.3), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
for i, f in enumerate(features):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'  {f}'
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(4)

# Trend summary
card = add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8))
add_text_box(slide, Inches(0.8), Inches(4.4), Inches(11), Inches(0.4),
             '2025-2026 Design Trends in SE Asia', font_size=20, bold=True, color=BLACK)

trends = [
    ('"Ugly-Cute" Aesthetic', 'LABUBU-style merging quirky + cute'),
    ('Candy/Pastel Colors', 'Soft, dreamy color palettes'),
    ('Y2K Revival', 'Millennium aesthetic comeback'),
    ('Cultural Fusion', 'Local artist collaborations'),
    ('Plush/Soft Materials', 'Plush keychain blind boxes trending'),
    ('Social-First Design', 'TikTok/Instagram-friendly visuals'),
]
for i, (trend, desc) in enumerate(trends):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(5) + row * Inches(1)
    add_text_box(slide, x, y, Inches(3.8), Inches(0.4),
                 trend, font_size=15, bold=True, color=DARK_PINK)
    add_text_box(slide, x, y + Inches(0.35), Inches(3.8), Inches(0.4),
                 desc, font_size=12, color=GRAY)

# ============================================================
# Slide 14: Project Deliverables Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             'Project Deliverables Summary', font_size=36, bold=True, color=WHITE)

deliverables = [
    ('Market Data Report', 'SE Asia 6-country market analysis\nJSON + Markdown formats\n2023-2025 data coverage', PINK),
    ('Web Templates', '10 HTML design concepts\nZero external dependencies\nInteractive prototypes', BLUE),
    ('Style Evaluation', '7-dimension scoring system\n4 brand analysis\n3 recommendation tiers', GREEN),
    ('Acceptance Report', 'Overall score: 8.2/10\n4 verification dimensions\nQuality issue tracking', YELLOW),
]
for i, (title, desc, clr) in enumerate(deliverables):
    x = Inches(0.3) + i * Inches(3.2)
    card = add_card(slide, x, Inches(1.3), Inches(3), Inches(3.5))
    # Icon circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(1), Inches(1.5), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = clr
    circle.line.fill.background()
    num = ['01', '02', '03', '04'][i]
    add_text_box(slide, x + Inches(1), Inches(1.6), Inches(1), Inches(0.8),
                 num, font_size=24, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.6), Inches(0.5),
                 title, font_size=18, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    txBox = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.2), Inches(2.6), Inches(1.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f'  {line}'
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER

# Acceptance scores
card = add_card(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8))
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(0.4),
             'Acceptance Scores', font_size=20, bold=True, color=WHITE)
score_items = [
    ('Data Quality', '9/10', '30%'), ('Web Templates', '7/10', '35%'),
    ('Style Analysis', '9/10', '25%'), ('Structure', '8/10', '10%'),
    ('TOTAL', '8.2/10', '100%'),
]
for i, (label, score, weight) in enumerate(score_items):
    x = Inches(0.8) + i * Inches(2.4)
    clr = YELLOW if label == 'TOTAL' else RGBColor(0xBB, 0xBB, 0xBB)
    bld = label == 'TOTAL'
    add_text_box(slide, x, Inches(5.8), Inches(2.2), Inches(0.3),
                 label, font_size=13, color=RGBColor(0x99, 0x99, 0x99),
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.1), Inches(2.2), Inches(0.5),
                 score, font_size=22, bold=bld, color=clr,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.5), Inches(2.2), Inches(0.3),
                 f'Weight: {weight}', font_size=11, color=RGBColor(0x77, 0x77, 0x77),
                 alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 15: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_PINK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             'Next Steps & Recommendations', font_size=36, bold=True, color=DARK_PINK)

steps = [
    ('Phase 1: Responsive Design', 'Add @media queries to all templates\nMobile-first approach\nTest on 320px, 768px, 1024px breakpoints', PINK),
    ('Phase 2: A/B Testing', 'Test top 3 templates in target markets\nMeasure conversion rates\nThailand, Vietnam, Philippines focus', BLUE),
    ('Phase 3: Local Integration', 'Integrate local payment gateways\nConnect TikTok Shop / Shopee\nAdd multi-language support', GREEN),
    ('Phase 4: Content & Assets', 'Replace emoji with real product images\nAdd IP character illustrations\nCreate box-opening animations', YELLOW),
]
for i, (phase, desc, clr) in enumerate(steps):
    x = Inches(0.5) + i * Inches(3.1)
    card = add_card(slide, x, Inches(1.2), Inches(2.9), Inches(4))
    # Phase number
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(1), Inches(1.4), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = clr
    circle.line.fill.background()
    add_text_box(slide, x + Inches(1), Inches(1.5), Inches(0.8), Inches(0.6),
                 str(i + 1), font_size=24, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.4), Inches(2.5), Inches(0.5),
                 phase, font_size=16, bold=True, color=BLACK,
                 alignment=PP_ALIGN.CENTER)
    txBox = slide.shapes.add_textbox(x + Inches(0.2), Inches(3), Inches(2.5), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f'  {line}'
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.font.name = 'Microsoft YaHei'
        p.space_before = Pt(6)

# Priority matrix
card = add_card(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5))
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.4),
             'Priority Matrix', font_size=18, bold=True, color=BLACK)
priorities = [
    ('HIGH', 'Responsive CSS for all templates', DARK_PINK),
    ('HIGH', 'Mobile-first optimization', DARK_PINK),
    ('MEDIUM', 'A/B testing in target markets', BLUE),
    ('MEDIUM', 'Local payment integration', BLUE),
    ('LOW', 'Replace emoji with real images', GRAY),
    ('LOW', 'JavaScript error handling', GRAY),
]
for i, (level, task, clr) in enumerate(priorities):
    x = Inches(0.8) + (i % 3) * Inches(4.1)
    y = Inches(6.1) + (i // 3) * Inches(0.4)
    add_text_box(slide, x, y, Inches(0.8), Inches(0.35),
                 f'[{level}]', font_size=12, bold=True, color=clr)
    add_text_box(slide, x + Inches(0.8), y, Inches(3), Inches(0.35),
                 task, font_size=12, color=GRAY)

# ============================================================
# Slide 16: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

# Decorative
for cx, cy, sz, clr in [
    (Inches(2), Inches(2), Inches(1.5), PINK),
    (Inches(10), Inches(5), Inches(2), BLUE),
    (Inches(11), Inches(1), Inches(1), GREEN),
    (Inches(1), Inches(6), Inches(0.8), YELLOW),
]:
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, sz, sz)
    c.fill.solid()
    c.fill.fore_color.rgb = clr
    c.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
             'Thank You', font_size=52, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4), Inches(10), Inches(0.6),
             'Southeast Asia Blind Box Market Web Design Proposal',
             font_size=20, color=PINK, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5), Inches(10), Inches(0.5),
             '2026.06', font_size=16, color=RGBColor(0x88, 0x88, 0x88),
             alignment=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = 'F:/测试工具/blindbox-project/presentation/SE-Asia-BlindBox-Market-Presentation.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
